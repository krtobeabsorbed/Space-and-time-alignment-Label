from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # 定义一个简单的辅助函数来设置文本格式
    def setup_text(shape, text, size=18, bold=False):
        tf = shape.text_frame
        tf.clear() # 清除默认格式
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = 'Arial' # 通用字体

    # 定义辅助函数：添加带标题的幻灯片
    def add_slide(title_text, content_text_list):
        slide_layout = prs.slide_layouts[1] # 1 是标题+内容布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = title_text
        
        # 设置内容
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for item in content_text_list:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.level = 0
            p.space_after = Pt(10)
            
        return slide

    # =======================================================
    # Slide 1: 封面
    # =======================================================
    slide_layout = prs.slide_layouts[0] # 0 是封面布局
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    
    title.text = "多模态数据采集中的雷达-视觉时空对齐方案"
    subtitle.text = "Spatial-Temporal Alignment for Radar-Vision Fusion Dataset\n\n汇报人：[你的名字]\n项目：小米智能家居多模态数据集"

    # =======================================================
    # Slide 2: 问题背景与挑战
    # =======================================================
    content = [
        "核心任务：",
        "将无语义的毫米波雷达点云（3D）与高分辨率图像（2D）进行像素级对齐。",
        "",
        "面临的三大挑战：",
        "1. 不可见性 (Invisibility)：雷达点稀疏且无纹理，无法提取特征点进行自动匹配。",
        "2. 数据退化 (Data Degeneracy)：室内强反射（墙壁、家具）导致90%以上的数据为静止杂波。",
        "3. 几何失真 (Geometric Distortion)：雷达与相机视角差异导致严重的非线性透视变换。"
    ]
    add_slide("问题背景与挑战 (Challenges)", content)

    # =======================================================
    # Slide 3: 技术路线图 (Methodology)
    # =======================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "技术路线图 (Methodology Pipeline)"
    
    # 这一页留白给流程图
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = "【此处请插入流程图】\nRaw Data -> Static Clutter Removal -> Coordinate Transform -> Interactive Tuning -> Fusion"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 0, 0) # 红色提示

    # =======================================================
    # Slide 4: 关键步骤一：静态杂波抑制
    # =======================================================
    content = [
        "问题诊断：",
        "原始雷达数据中，墙壁和家具的反射信号强度远高于人体，导致跟踪算法'死锁'。",
        "",
        "解决方案 (Static Clutter Removal)：",
        "- 构建空间直方图 (Spatial Histogram)，识别高频驻留点。",
        "- 发现并剔除死锁坐标 (如: X≈0.24, Y≈4.72)。",
        "- 结果：剔除 90.3% 的静止帧，仅保留动态人体轨迹，信噪比显著提升。"
    ]
    slide = add_slide("关键步骤一：静态杂波抑制", content)
    
    # 红色提示框
    textbox = slide.shapes.add_textbox(Inches(5), Inches(5), Inches(4), Inches(1))
    p = textbox.text_frame.paragraphs[0]
    p.text = "【此处插入 clean_radar_track.py 生成的 Before/After 对比图】"
    p.font.color.rgb = RGBColor(255, 0, 0)

    # =======================================================
    # Slide 5: 关键步骤二：数学模型
    # =======================================================
    content = [
        "坐标变换模型 (针孔相机模型)：",
        "将雷达坐标 (Xr, Yr, Zr) 映射到 图像坐标 (u, v)",
        "",
        "s * [u, v, 1]^T = K * [R | T] * [Xr, Yr, Zr, 1]^T",
        "",
        "- K: 相机内参 (焦距 f=4mm, 传感器尺寸 5.9mm)",
        "- R: 旋转矩阵 (修正俯拍 Pitch 和偏航 Yaw 角度)",
        "- T: 平移向量 (修正雷达与相机的物理位移)",
        "- Mirror Correction: 针对传感器朝向进行 X 轴镜像处理"
    ]
    add_slide("关键步骤二：数学模型 (Mathematical Model)", content)

    # =======================================================
    # Slide 6: 关键步骤三：交互式时空对齐
    # =======================================================
    content = [
        "创新点：Human-in-the-loop 交互式标定",
        "针对雷达数据稀疏、PnP 算法易发散的问题，开发实时反馈调参工具。",
        "",
        "最终标定参数 (Final Parameters)：",
        "1. Time Offset (时间偏移): 0.0s (数据清洗后对齐)",
        "2. Translation (空间位移): T = [1.7, 4.5, 0.5] (米)",
        "   *注：Ty=4.5m 包含了对广角畸变和透视误差的等效补偿",
        "3. Rotation (姿态修正): Pitch = 25.0° (俯拍)",
        "4. Mirror X: True (开启镜像)"
    ]
    add_slide("关键步骤三：交互式时空对齐", content)

    # =======================================================
    # Slide 7: 成果展示
    # =======================================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "成果展示 (Demo & Results)"
    
    textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    p = textbox.text_frame.paragraphs[0]
    p.text = "【此处插入 output_fusion_final.mp4 的视频截图】\n展示红点紧密跟随人体脚底的效果"
    p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(255, 0, 0)

    # =======================================================
    # Slide 8: 结论
    # =======================================================
    content = [
        "1. 成功实现了 Radar 与 Camera 的亚秒级时间同步与像素级空间对齐。",
        "2. 解决了室内强静态反射干扰的问题，提取出纯净的人体轨迹。",
        "3. 生成了包含 (Pixel_U, Pixel_V, Depth_Z) 的高质量 CSV 数据集。",
        "4. 该方案参数已固定，可直接批量应用于后续的大规模数据采集。"
    ]
    add_slide("结论 (Conclusion)", content)

    # 保存
    output_file = 'Xiaomi_Report_Alignment.pptx'
    prs.save(output_file)
    print(f"✅ PPT 已成功生成: {output_file}")
    print("请打开 PPT，并将红字部分的占位符替换为真实的图片。")

if __name__ == "__main__":
    create_presentation()